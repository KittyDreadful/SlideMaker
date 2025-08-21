from pptx import Presentation
from config import TEMPLATE_PATH

def load_template():
    return Presentation(TEMPLATE_PATH)

def duplicate_slide(prs, slide):
    layout = prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.has_text_frame:
            new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_shape.text = shape.text
    return new_slide

def replace_placeholders(slide, row):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text
            for key in row.keys():
                placeholder = f"{{{{{key}}}}}"
                if placeholder in text:
                    text = text.replace(placeholder, str(row[key]))
            shape.text = text
